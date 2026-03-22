"""
Generate IT News Platform presentation (.pptx)
"""

import io
import os
import textwrap

import qrcode
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

# ── Colour palette ──────────────────────────────────────────────────────────
BG_DARK   = RGBColor(0x0D, 0x11, 0x17)   # deep navy
BG_CARD   = RGBColor(0x16, 0x1B, 0x22)   # card background
ACCENT    = RGBColor(0x58, 0xA6, 0xFF)   # bright blue
ACCENT2   = RGBColor(0x3F, 0xB9, 0x50)   # green
ACCENT3   = RGBColor(0xD2, 0xA8, 0xFF)   # purple
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY      = RGBColor(0x8B, 0x94, 0x9E)
LIGHT     = RGBColor(0xC9, 0xD1, 0xD9)
ORANGE    = RGBColor(0xF0, 0x88, 0x3E)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

OUT_DIR = os.path.dirname(os.path.abspath(__file__))


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline(slide, left, top, width, height, lines, font_size=16,
                  color=WHITE, bold=False, spacing=1.2, font_name="Segoe UI",
                  alignment=PP_ALIGN.LEFT):
    """lines: list of (text, color|None) or just str"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, tuple):
            txt, clr = item
        else:
            txt, clr = item, color
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(font_size)
        p.font.color.rgb = clr or color
        p.font.bold = bold
        p.font.name = font_name
        p.space_after = Pt(font_size * (spacing - 1) + 2)
        p.alignment = alignment
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color, text="",
                     font_size=14, font_color=WHITE, bold=False):
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.font.name = "Segoe UI"
    shape.text_frame.paragraphs[0].space_before = Pt(0)
    shape.text_frame.paragraphs[0].space_after = Pt(0)
    try:
        shape.text_frame.margin_top = Emu(0)
        shape.text_frame.margin_bottom = Emu(0)
    except Exception:
        pass
    return shape


def make_qr(url: str, size: int = 200) -> io.BytesIO:
    qr = qrcode.QRCode(version=1, box_size=8, border=2,
                        error_correction=qrcode.constants.ERROR_CORRECT_M)
    qr.add_data(url)
    qr.make(fit=True)
    img = qr.make_image(fill_color="white", back_color="#0D1117").convert("RGBA")
    img = img.resize((size, size), Image.LANCZOS)
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


def make_icon(emoji: str, size: int = 120, bg_color="#161B22") -> io.BytesIO:
    img = Image.new("RGBA", (size, size), bg_color)
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


def add_arrow(slide, x1, y1, x2, y2, color=ACCENT):
    """Draw arrow between two points"""
    from pptx.enum.shapes import MSO_SHAPE
    connector = slide.shapes.add_connector(
        1, x1, y1, x2, y2  # MSO_CONNECTOR_TYPE.STRAIGHT
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(2)


def create_gradient_image(width, height, color1, color2) -> io.BytesIO:
    """Create a simple gradient image."""
    img = Image.new("RGB", (width, height))
    draw = ImageDraw.Draw(img)
    r1, g1, b1 = int(color1[1:3], 16), int(color1[3:5], 16), int(color1[5:7], 16)
    r2, g2, b2 = int(color2[1:3], 16), int(color2[3:5], 16), int(color2[5:7], 16)
    for y in range(height):
        ratio = y / height
        r = int(r1 + (r2 - r1) * ratio)
        g = int(g1 + (g2 - g1) * ratio)
        b = int(b1 + (b2 - b1) * ratio)
        draw.line([(0, y), (width, y)], fill=(r, g, b))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


def create_tech_pattern(width, height) -> io.BytesIO:
    """Create a subtle tech-inspired background pattern."""
    img = Image.new("RGBA", (width, height), (13, 17, 23, 255))
    draw = ImageDraw.Draw(img)
    # Grid dots
    for x in range(0, width, 40):
        for y in range(0, height, 40):
            draw.ellipse([x-1, y-1, x+1, y+1], fill=(88, 166, 255, 25))
    # Subtle lines
    for x in range(0, width, 120):
        draw.line([(x, 0), (x, height)], fill=(88, 166, 255, 12), width=1)
    for y in range(0, height, 120):
        draw.line([(0, y), (width, y)], fill=(88, 166, 255, 12), width=1)
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


def create_node_graph_image(width, height) -> io.BytesIO:
    """Create an abstract knowledge graph visualization."""
    img = Image.new("RGBA", (width, height), (0, 0, 0, 0))
    draw = ImageDraw.Draw(img)
    import random
    random.seed(42)
    nodes = []
    for _ in range(12):
        x = random.randint(40, width - 40)
        y = random.randint(40, height - 40)
        nodes.append((x, y))
    # Draw edges
    for i in range(len(nodes)):
        for j in range(i + 1, len(nodes)):
            dx = nodes[j][0] - nodes[i][0]
            dy = nodes[j][1] - nodes[i][1]
            dist = (dx ** 2 + dy ** 2) ** 0.5
            if dist < 200:
                alpha = max(10, int(80 - dist / 3))
                draw.line([nodes[i], nodes[j]], fill=(88, 166, 255, alpha), width=1)
    # Draw nodes
    colors = [(88, 166, 255), (63, 185, 80), (210, 168, 255), (240, 136, 62)]
    for i, (x, y) in enumerate(nodes):
        c = colors[i % len(colors)]
        r = random.randint(6, 14)
        draw.ellipse([x - r, y - r, x + r, y + r], fill=(*c, 180))
        draw.ellipse([x - r - 3, y - r - 3, x + r + 3, y + r + 3], outline=(*c, 60), width=1)
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


def create_pipeline_arrow_image(width=120, height=40) -> io.BytesIO:
    """Arrow image for pipeline."""
    img = Image.new("RGBA", (width, height), (0, 0, 0, 0))
    draw = ImageDraw.Draw(img)
    # Arrow body
    draw.rectangle([10, height // 2 - 3, width - 30, height // 2 + 3], fill=(88, 166, 255, 180))
    # Arrow head
    draw.polygon([(width - 30, height // 2 - 12), (width - 5, height // 2), (width - 30, height // 2 + 12)],
                 fill=(88, 166, 255, 180))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


def build_presentation():
    prs = Presentation()
    prs.slide_width = Emu(12192000)  # 13.333 inches
    prs.slide_height = Emu(6858000)  # 7.5 inches

    # ═══════════════════════════════════════════════════════════════
    # SLIDE 1: Title
    # ═══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BG_DARK)

    # Background pattern
    pattern = create_tech_pattern(1920, 1080)
    slide.shapes.add_picture(pattern, 0, 0, Emu(12192000), Emu(6858000))

    # Graph decoration
    graph_img = create_node_graph_image(600, 500)
    slide.shapes.add_picture(graph_img, Inches(8.5), Inches(1.5), Inches(4.5), Inches(3.8))

    # Accent line top
    add_rounded_rect(slide, Inches(1), Inches(1.3), Inches(0.6), Pt(4), ACCENT)

    add_textbox(slide, Inches(1), Inches(1.5), Inches(8), Inches(1.2),
                "IT News Platform", font_size=48, color=WHITE, bold=True)
    add_textbox(slide, Inches(1), Inches(2.5), Inches(8), Inches(0.8),
                "Интеллектуальная платформа агрегации\nIT-новостей с графом знаний и RAG",
                font_size=22, color=LIGHT)

    # Tech badges
    techs = ["FastAPI", "React + TS", "Neo4j", "Qdrant", "Groq LLM"]
    for i, tech in enumerate(techs):
        add_rounded_rect(slide, Inches(1 + i * 1.7), Inches(3.8), Inches(1.5), Inches(0.45),
                         BG_CARD, tech, font_size=13, font_color=ACCENT, bold=True)

    add_textbox(slide, Inches(1), Inches(5.5), Inches(6), Inches(0.5),
                "Хакатон 2026  •  [Название команды]",
                font_size=16, color=GRAY)

    # ═══════════════════════════════════════════════════════════════
    # SLIDE 2: Problem & Solution
    # ═══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    pattern.seek(0)
    slide.shapes.add_picture(pattern, 0, 0, Emu(12192000), Emu(6858000))

    add_rounded_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT)
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.7),
                "Проблема и Решение", font_size=36, color=WHITE, bold=True)

    # Problem card
    add_rounded_rect(slide, Inches(0.6), Inches(1.4), Inches(5.6), Inches(5.2), BG_CARD)
    add_textbox(slide, Inches(1.0), Inches(1.6), Inches(5), Inches(0.5),
                "⚡  Проблема", font_size=24, color=ORANGE, bold=True)
    add_multiline(slide, Inches(1.0), Inches(2.3), Inches(5), Inches(4), [
        ("•  Тысячи IT-новостей ежедневно —\n   невозможно отследить всё вручную", LIGHT),
        ("•  Разрозненные источники без\n   единой точки доступа", LIGHT),
        ("•  Нет связей между событиями,\n   компаниями и персонами", LIGHT),
        ("•  Простой keyword-поиск не понимает\n   смысл запроса", LIGHT),
    ], font_size=17, spacing=1.5)

    # Solution card
    add_rounded_rect(slide, Inches(6.8), Inches(1.4), Inches(5.9), Inches(5.2), BG_CARD)
    add_textbox(slide, Inches(7.2), Inches(1.6), Inches(5), Inches(0.5),
                "✅  Наше решение", font_size=24, color=ACCENT2, bold=True)
    add_multiline(slide, Inches(7.2), Inches(2.3), Inches(5.2), Inches(4), [
        ("•  Автоматический сбор новостей\n   из RSS и NewsAPI", LIGHT),
        ("•  Граф знаний (Neo4j) — сущности\n   и связи между ними", LIGHT),
        ("•  Семантический поиск на\n   эмбеддингах (Qdrant)", LIGHT),
        ("•  RAG с LLM — ответы на вопросы\n   по актуальным новостям", LIGHT),
    ], font_size=17, spacing=1.5)

    # ═══════════════════════════════════════════════════════════════
    # SLIDE 3: Architecture (clean vertical flow)
    # ═══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    pattern.seek(0)
    slide.shapes.add_picture(pattern, 0, 0, Emu(12192000), Emu(6858000))

    add_rounded_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT3)
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.7),
                "Архитектура платформы", font_size=36, color=WHITE, bold=True)

    # ── Row 1: Sources ──────────────────────────────────────────
    row1_y = 1.35
    add_textbox(slide, Inches(0.6), Inches(row1_y), Inches(1.5), Inches(0.35),
                "ИСТОЧНИКИ", font_size=11, color=GRAY, bold=True)
    sources = ["TechCrunch RSS", "Wired RSS", "Ars Technica RSS", "NewsAPI"]
    for i, src in enumerate(sources):
        add_rounded_rect(slide, Inches(2.2 + i * 2.6), Inches(row1_y - 0.05),
                         Inches(2.3), Inches(0.5), BG_CARD, src,
                         font_size=12, font_color=LIGHT)

    # Arrow ▼
    add_textbox(slide, Inches(5.8), Inches(row1_y + 0.45), Inches(1), Inches(0.35),
                "▼", font_size=20, color=ACCENT, alignment=PP_ALIGN.CENTER)

    # ── Row 2: Ingestion ────────────────────────────────────────
    row2_y = 2.15
    add_textbox(slide, Inches(0.6), Inches(row2_y + 0.05), Inches(1.5), Inches(0.35),
                "INGESTION", font_size=11, color=GRAY, bold=True)
    add_rounded_rect(slide, Inches(2.2), Inches(row2_y), Inches(10.5), Inches(0.55),
                     RGBColor(0x14, 0x25, 0x35),
                     "RSS / NewsAPI  →  очистка HTML  →  нормализация  →  language detect  →  дедупликация по URL  →  PostgreSQL",
                     font_size=13, font_color=ACCENT)

    # Arrow ▼
    add_textbox(slide, Inches(5.8), Inches(row2_y + 0.55), Inches(1), Inches(0.35),
                "▼", font_size=20, color=ACCENT, alignment=PP_ALIGN.CENTER)

    # ── Row 3: Storage layer (4 DBs side by side) ──────────────
    row3_y = 3.0
    add_textbox(slide, Inches(0.6), Inches(row3_y + 0.15), Inches(1.5), Inches(0.35),
                "ХРАНИЛИЩА", font_size=11, color=GRAY, bold=True)

    db_items = [
        ("PostgreSQL", "Статьи, метаданные\nиндексы, embedding_status", ACCENT2, RGBColor(0x14, 0x28, 0x14)),
        ("Neo4j", "Граф знаний\nArticle → Entity → Relations", ACCENT, RGBColor(0x14, 0x18, 0x30)),
        ("Qdrant", "Векторное хранилище\n384d эмбеддинги, cosine", ORANGE, RGBColor(0x30, 0x1C, 0x10)),
        ("Redis", "Кэш и очереди\nинфраструктурный слой", RGBColor(0xFF, 0x60, 0x60), RGBColor(0x2A, 0x12, 0x12)),
    ]
    for i, (name, desc, color, bg) in enumerate(db_items):
        x = Inches(2.2 + i * 2.6)
        add_rounded_rect(slide, x, Inches(row3_y), Inches(2.3), Inches(1.0), bg)
        add_textbox(slide, x + Inches(0.15), Inches(row3_y + 0.05), Inches(2), Inches(0.35),
                    name, font_size=15, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.1), Inches(row3_y + 0.4), Inches(2.1), Inches(0.55),
                    desc, font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)

    # Arrow ▼
    add_textbox(slide, Inches(5.8), Inches(row3_y + 1.0), Inches(1), Inches(0.35),
                "▼", font_size=20, color=ACCENT, alignment=PP_ALIGN.CENTER)

    # ── Row 4: Processing (NLP + Embeddings + Clustering) ──────
    row4_y = 4.35
    add_textbox(slide, Inches(0.6), Inches(row4_y + 0.1), Inches(1.5), Inches(0.35),
                "ОБРАБОТКА", font_size=11, color=GRAY, bold=True)

    proc_items = [
        ("NLP Pipeline", "Rule-based NER\nPERSON • ORG • LOCATION\n→ Neo4j граф", ACCENT3, RGBColor(0x22, 0x14, 0x28)),
        ("Embedding Pipeline", "MiniLM-L12-v2\nL2-нормализация\n→ Qdrant upsert", ORANGE, RGBColor(0x28, 0x1E, 0x10)),
        ("Кластеризация", "HDBSCAN (≥50)\nKMeans (<50)\ncluster_id → PG + Qdrant", ACCENT2, RGBColor(0x14, 0x24, 0x14)),
    ]
    proc_w = 3.3
    for i, (name, desc, color, bg) in enumerate(proc_items):
        x = Inches(2.2 + i * (proc_w + 0.35))
        add_rounded_rect(slide, x, Inches(row4_y), Inches(proc_w), Inches(1.1), bg)
        add_textbox(slide, x + Inches(0.15), Inches(row4_y + 0.05), Inches(proc_w - 0.3), Inches(0.35),
                    name, font_size=14, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.1), Inches(row4_y + 0.38), Inches(proc_w - 0.2), Inches(0.7),
                    desc, font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)

    # Arrow ▼
    add_textbox(slide, Inches(5.8), Inches(row4_y + 1.1), Inches(1), Inches(0.35),
                "▼", font_size=20, color=ACCENT, alignment=PP_ALIGN.CENTER)

    # ── Row 5: RAG (full width) ────────────────────────────────
    row5_y = 5.8
    add_textbox(slide, Inches(0.6), Inches(row5_y + 0.15), Inches(1.5), Inches(0.35),
                "RAG", font_size=11, color=GRAY, bold=True)
    add_rounded_rect(slide, Inches(2.2), Inches(row5_y), Inches(10.5), Inches(0.6),
                     RGBColor(0x18, 0x1E, 0x2E),
                     "Запрос → Query Embedding → Qdrant top-K → Neo4j (entities + edges) → Groq LLM (3-model cascade) → Ответ + Источники",
                     font_size=13, font_color=WHITE, bold=True)
    add_multiline(slide, Inches(2.2), Inches(row5_y + 0.6), Inches(10.5), Inches(0.35), [
        ("Fallback:  gpt-oss-120b  →  llama-3.3-70b  →  llama-3.1-8b  →  retrieval-only", GRAY),
    ], font_size=11, alignment=PP_ALIGN.CENTER)

    # ── Footer: Frontend + Deploy ───────────────────────────────
    add_textbox(slide, Inches(0.6), Inches(row5_y + 1.0), Inches(12), Inches(0.35),
                "🖥️ Frontend: React + TypeScript + D3.js    |    🐳 Docker Compose    |    ☁️ Railway / Vercel",
                font_size=12, color=GRAY)

    # ═══════════════════════════════════════════════════════════════
    # SLIDE 4: Data Pipeline Deep Dive
    # ═══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    pattern.seek(0)
    slide.shapes.add_picture(pattern, 0, 0, Emu(12192000), Emu(6858000))

    add_rounded_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT2)
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.7),
                "Конвейер обработки данных", font_size=36, color=WHITE, bold=True)

    # Pipeline steps as cards
    steps = [
        ("1", "Ingestion", "RSS / NewsAPI\nасинхронный сбор\nhttpx + feedparser", ACCENT),
        ("2", "Очистка", "HTML → чистый текст\nBeautifulSoup\nnormalize + langdetect", ACCENT2),
        ("3", "Хранение", "PostgreSQL\nдедупликация по URL\nиндексы по source", ACCENT3),
        ("4", "Embeddings", "MiniLM-L12-v2\n384-d вектора\nL2 нормализация", ORANGE),
        ("5", "Индексация", "Qdrant upsert\npayload: source, date\nentity_names, cluster", RGBColor(0xFF, 0x60, 0x60)),
    ]

    for i, (num, title, desc, color) in enumerate(steps):
        x = Inches(0.5 + i * 2.55)
        y = Inches(1.6)
        add_rounded_rect(slide, x, y, Inches(2.3), Inches(2.8), BG_CARD)
        # Number circle
        add_rounded_rect(slide, x + Inches(0.1), y + Inches(0.15), Inches(0.45), Inches(0.45),
                         color, num, font_size=18, font_color=BG_DARK, bold=True)
        add_textbox(slide, x + Inches(0.65), y + Inches(0.15), Inches(1.5), Inches(0.45),
                    title, font_size=18, color=color, bold=True)
        add_multiline(slide, x + Inches(0.15), y + Inches(0.8), Inches(2), Inches(1.8), [
            (desc, LIGHT),
        ], font_size=14, spacing=1.3)

        # Arrow between
        if i < len(steps) - 1:
            add_textbox(slide, x + Inches(2.3), Inches(2.6), Inches(0.3), Inches(0.4),
                        "→", font_size=20, color=GRAY)

    # Bottom section: Clustering
    add_rounded_rect(slide, Inches(0.5), Inches(4.8), Inches(6), Inches(1.8), BG_CARD)
    add_textbox(slide, Inches(0.9), Inches(4.95), Inches(5), Inches(0.5),
                "🔬  Кластеризация", font_size=20, color=ACCENT3, bold=True)
    add_multiline(slide, Inches(0.9), Inches(5.5), Inches(5.2), Inches(1.2), [
        ("•  < 50 статей → KMeans (адаптивное k)", LIGHT),
        ("•  ≥ 50 статей → HDBSCAN (плотностная)", LIGHT),
        ("•  cluster_id синхронизируется в PostgreSQL и Qdrant", LIGHT),
        ("•  Автоматический recluster каждые 30 мин", LIGHT),
    ], font_size=14, spacing=1.4)

    # NER/Graph section
    add_rounded_rect(slide, Inches(6.9), Inches(4.8), Inches(5.9), Inches(1.8), BG_CARD)
    add_textbox(slide, Inches(7.3), Inches(4.95), Inches(5), Inches(0.5),
                "🕸️  NLP → Граф знаний", font_size=20, color=ACCENT, bold=True)
    add_multiline(slide, Inches(7.3), Inches(5.5), Inches(5.2), Inches(1.2), [
        ("•  Rule-based NER: PERSON, ORG, LOCATION", LIGHT),
        ("•  Relation extraction: ASSOCIATED_WITH, LOCATED_IN", LIGHT),
        ("•  Neo4j: Article → MENTIONS → Entity", LIGHT),
        ("•  Entity → RELATED_TO → Entity (с article_ids)", LIGHT),
    ], font_size=14, spacing=1.4)

    # ═══════════════════════════════════════════════════════════════
    # SLIDE 5: RAG Pipeline
    # ═══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    pattern.seek(0)
    slide.shapes.add_picture(pattern, 0, 0, Emu(12192000), Emu(6858000))

    add_rounded_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ORANGE)
    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(12), Inches(0.7),
                "RAG: Retrieval-Augmented Generation", font_size=36, color=WHITE, bold=True)

    # RAG flow steps
    rag_steps = [
        ("Вопрос\nпользователя", ACCENT, "POST /api/search\n{question, top_k,\nuse_graph, filters}"),
        ("Query\nEmbedding", ACCENT2, "MiniLM-L12-v2\nтот же трансформер\nчто и для статей"),
        ("Qdrant\nRetrieval", ORANGE, "Cosine similarity\ntop-K × 3 → rerank\nfreshness boost +0.05"),
        ("Neo4j\nGraph", ACCENT3, "entities + edges\nпо найденным статьям\nmax 15 entities, 20 rels"),
        ("LLM\nGeneration", RGBColor(0xFF, 0x60, 0x60), "Groq: 3-model cascade\nstrict JSON output\nconfidence score"),
        ("Ответ +\nИсточники", WHITE, "answer + sources +\nentities + graph_edges\n+ retrieval_debug"),
    ]

    for i, (title, color, desc) in enumerate(rag_steps):
        x = Inches(0.3 + i * 2.15)
        y = Inches(1.5)
        add_rounded_rect(slide, x, y, Inches(1.95), Inches(1.0), BG_CARD)
        add_textbox(slide, x + Inches(0.1), y + Inches(0.05), Inches(1.75), Inches(0.9),
                    title, font_size=15, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_multiline(slide, x + Inches(0.05), y + Inches(1.1), Inches(1.85), Inches(1.2), [
            (desc, GRAY),
        ], font_size=11, alignment=PP_ALIGN.CENTER, spacing=1.2)

        if i < len(rag_steps) - 1:
            add_textbox(slide, x + Inches(1.95), Inches(1.75), Inches(0.25), Inches(0.4),
                        "→", font_size=18, color=ACCENT)

    # Degradation strategy
    add_rounded_rect(slide, Inches(0.5), Inches(4.2), Inches(12.3), Inches(2.5), BG_CARD)
    add_textbox(slide, Inches(0.9), Inches(4.35), Inches(8), Inches(0.5),
                "🛡️  Стратегия отказоустойчивости (Graceful Degradation)", font_size=20, color=ACCENT, bold=True)

    # Three fallback levels
    levels = [
        ("Уровень 1", "gpt-oss-120b", "Основная модель\n120B параметров", ACCENT2),
        ("Уровень 2", "llama-3.3-70b", "Fallback модель\n70B параметров", ORANGE),
        ("Уровень 3", "llama-3.1-8b", "Быстрая модель\n8B параметров", RGBColor(0xFF, 0x60, 0x60)),
        ("Уровень 4", "Retrieval-only", "Без LLM\nТолько статьи + граф", GRAY),
    ]

    for i, (level, model, desc, color) in enumerate(levels):
        x = Inches(0.7 + i * 3.1)
        add_rounded_rect(slide, x, Inches(5.0), Inches(2.8), Inches(1.4),
                         RGBColor(0x1A, 0x1F, 0x28))
        add_textbox(slide, x + Inches(0.15), Inches(5.05), Inches(2.5), Inches(0.35),
                    level, font_size=12, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.15), Inches(5.35), Inches(2.5), Inches(0.35),
                    model, font_size=15, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_multiline(slide, x + Inches(0.15), Inches(5.75), Inches(2.5), Inches(0.6), [
            (desc, GRAY),
        ], font_size=12, alignment=PP_ALIGN.CENTER, spacing=1.2)

        if i < len(levels) - 1:
            add_textbox(slide, x + Inches(2.8), Inches(5.35), Inches(0.35), Inches(0.4),
                        "→", font_size=18, color=color)

    # ═══════════════════════════════════════════════════════════════
    # SLIDE 6: Frontend & Визуализация
    # ═══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    pattern.seek(0)
    slide.shapes.add_picture(pattern, 0, 0, Emu(12192000), Emu(6858000))

    CYAN = RGBColor(0x79, 0xC0, 0xFF)
    add_rounded_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), CYAN)
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.7),
                "Frontend и визуализация", font_size=36, color=WHITE, bold=True)

    # ── Left: UI mockup layout ──────────────────────────────────
    add_rounded_rect(slide, Inches(0.5), Inches(1.4), Inches(6.0), Inches(5.5), BG_CARD)
    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5), Inches(0.4),
                "🖥️  Интерфейс платформы", font_size=18, color=CYAN, bold=True)

    # Search / Chat block
    add_rounded_rect(slide, Inches(0.8), Inches(2.1), Inches(2.5), Inches(2.0),
                     RGBColor(0x1A, 0x22, 0x30))
    add_textbox(slide, Inches(0.95), Inches(2.15), Inches(2.2), Inches(0.35),
                "🔍 RAG-чат", font_size=14, color=ACCENT, bold=True)
    add_multiline(slide, Inches(0.95), Inches(2.55), Inches(2.2), Inches(1.5), [
        ("Вопрос → ответ с источниками", GRAY),
        ("История запросов", GRAY),
        ("Фильтры: дата, тема,", GRAY),
        ("источник, язык", GRAY),
    ], font_size=11, spacing=1.3)

    # Graph block
    add_rounded_rect(slide, Inches(3.5), Inches(2.1), Inches(2.8), Inches(2.0),
                     RGBColor(0x1A, 0x22, 0x30))
    add_textbox(slide, Inches(3.65), Inches(2.15), Inches(2.5), Inches(0.35),
                "🕸️ Интерактивный граф", font_size=14, color=ACCENT3, bold=True)
    # Mini graph visual
    graph_mini = create_node_graph_image(350, 200)
    slide.shapes.add_picture(graph_mini, Inches(3.7), Inches(2.6), Inches(2.4), Inches(1.4))

    # Articles list block
    add_rounded_rect(slide, Inches(0.8), Inches(4.3), Inches(5.5), Inches(1.2),
                     RGBColor(0x1A, 0x22, 0x30))
    add_textbox(slide, Inches(0.95), Inches(4.35), Inches(5), Inches(0.35),
                "📰 Лента статей + кластеры", font_size=14, color=ACCENT2, bold=True)
    add_multiline(slide, Inches(0.95), Inches(4.75), Inches(5.2), Inches(0.7), [
        ("Статьи с подсветкой сущностей  •  Группировка по кластерам  •  Сортировка по дате / релевантности", GRAY),
    ], font_size=11)

    # Tech badge row
    front_techs = ["React", "TypeScript", "D3.js", "Material UI", "Zustand"]
    for i, t in enumerate(front_techs):
        add_rounded_rect(slide, Inches(0.7 + i * 1.15), Inches(5.7), Inches(1.05), Inches(0.35),
                         RGBColor(0x1A, 0x28, 0x3A), t, font_size=10, font_color=CYAN)

    # ── Right: Features & interactions ──────────────────────────
    add_rounded_rect(slide, Inches(6.8), Inches(1.4), Inches(6.0), Inches(2.6), BG_CARD)
    add_textbox(slide, Inches(7.1), Inches(1.5), Inches(5), Inches(0.4),
                "⚡  Интерактивность графа", font_size=18, color=ACCENT3, bold=True)
    add_multiline(slide, Inches(7.1), Inches(2.0), Inches(5.5), Inches(1.8), [
        ("•  Клик на узел → связанные статьи и сущности", LIGHT),
        ("•  Клик на связь → тип отношения и источники", LIGHT),
        ("•  Фильтр по времени → скрыть старые связи", LIGHT),
        ("•  Поиск по людям / компаниям → подсветка узлов", LIGHT),
        ("•  Drag & drop → перемещение узлов", LIGHT),
    ], font_size=14, spacing=1.3)

    # Deploy info
    add_rounded_rect(slide, Inches(6.8), Inches(4.2), Inches(6.0), Inches(2.4), BG_CARD)
    add_textbox(slide, Inches(7.1), Inches(4.3), Inches(5), Inches(0.4),
                "☁️  Деплой", font_size=18, color=ACCENT2, bold=True)
    add_multiline(slide, Inches(7.1), Inches(4.8), Inches(5.5), Inches(1.7), [
        ("•  Backend: Railway / Render", LIGHT),
        ("•  Frontend: Vercel / Netlify", LIGHT),
        ("•  БД: Neo4j Aura, Supabase (PostgreSQL)", LIGHT),
        ("•  Векторная БД: Qdrant Cloud", LIGHT),
        ("•  Локально: docker compose up --build", LIGHT),
    ], font_size=14, spacing=1.3)

    # ═══════════════════════════════════════════════════════════════
    # SLIDE 6: Tech Stack
    # ═══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    pattern.seek(0)
    slide.shapes.add_picture(pattern, 0, 0, Emu(12192000), Emu(6858000))

    add_rounded_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT2)
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.7),
                "Технологический стек", font_size=36, color=WHITE, bold=True)

    categories = [
        ("Backend", [
            ("FastAPI", "Async REST API"),
            ("SQLAlchemy 2.0", "Async ORM"),
            ("APScheduler", "Периодические задачи"),
            ("Pydantic v2", "Валидация данных"),
        ], ACCENT),
        ("Frontend", [
            ("React + TS", "SPA-приложение"),
            ("D3.js", "Визуализация графа"),
            ("Material UI", "UI-компоненты"),
            ("Zustand", "State management"),
        ], RGBColor(0x79, 0xC0, 0xFF)),
        ("Хранилища", [
            ("PostgreSQL 16", "Статьи и метаданные"),
            ("Neo4j 5.26", "Граф знаний"),
            ("Qdrant 1.13", "Векторный поиск"),
            ("Redis 7", "Кэш и очереди"),
        ], ACCENT2),
        ("ML / NLP", [
            ("MiniLM-L12-v2", "Мультиязычные эмбеддинги"),
            ("HDBSCAN", "Плотностная кластеризация"),
            ("KMeans", "Fallback кластеризация"),
            ("Rule-based NER", "Извлечение сущностей"),
        ], ACCENT3),
        ("LLM / RAG", [
            ("Groq API", "Инференс-провайдер"),
            ("gpt-oss-120b", "Основная LLM"),
            ("LLaMA 3.3 70B", "Fallback LLM"),
            ("LLaMA 3.1 8B", "Fast LLM"),
        ], ORANGE),
        ("Деплой", [
            ("Docker Compose", "Оркестрация"),
            ("Railway", "Backend хостинг"),
            ("Vercel", "Frontend хостинг"),
            ("Neo4j Aura", "Cloud граф"),
        ], RGBColor(0xFF, 0x60, 0x60)),
    ]

    for col, (cat_name, items, color) in enumerate(categories):
        x = Inches(0.2 + col * 2.2)
        add_rounded_rect(slide, x, Inches(1.5), Inches(2.0), Inches(5.2), BG_CARD)
        add_textbox(slide, x + Inches(0.1), Inches(1.6), Inches(1.8), Inches(0.45),
                    cat_name, font_size=15, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        # Accent line under category name
        add_rounded_rect(slide, x + Inches(0.2), Inches(2.05), Inches(1.6), Pt(2), color)

        for i, (tech, desc) in enumerate(items):
            y = Inches(2.25 + i * 1.1)
            add_textbox(slide, x + Inches(0.15), y, Inches(1.7), Inches(0.35),
                        tech, font_size=13, color=WHITE, bold=True)
            add_textbox(slide, x + Inches(0.15), y + Inches(0.3), Inches(1.7), Inches(0.35),
                        desc, font_size=10, color=GRAY)

    # ═══════════════════════════════════════════════════════════════
    # SLIDE 7: API & Features
    # ═══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    pattern.seek(0)
    slide.shapes.add_picture(pattern, 0, 0, Emu(12192000), Emu(6858000))

    add_rounded_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT)
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.7),
                "API эндпоинты и возможности", font_size=36, color=WHITE, bold=True)

    endpoints = [
        ("GET", "/health", "Healthcheck всех компонентов", ACCENT2),
        ("POST", "/ingestion/run", "Запуск сбора новостей", ACCENT),
        ("POST", "/indexing/run", "Генерация эмбеддингов и индексация", ACCENT),
        ("POST", "/clustering/run", "Перекластеризация статей", ACCENT3),
        ("GET", "/articles", "Список последних статей", ACCENT),
        ("POST", "/articles/{id}/graph", "Извлечение графа для статьи", ACCENT3),
        ("GET", "/api/search/semantic", "Семантический поиск по статьям", ORANGE),
        ("GET", "/api/clusters", "Список кластеров с примерами", ACCENT3),
        ("POST", "/api/search", "RAG: ответ на вопрос по новостям", RGBColor(0xFF, 0x60, 0x60)),
    ]

    for i, (method, path, desc, color) in enumerate(endpoints):
        y = Inches(1.4 + i * 0.62)
        # Method badge
        add_rounded_rect(slide, Inches(0.6), y, Inches(1.0), Inches(0.45),
                         color, method, font_size=13, font_color=BG_DARK, bold=True)
        # Path
        add_textbox(slide, Inches(1.8), y + Inches(0.02), Inches(4), Inches(0.42),
                    path, font_size=15, color=WHITE, bold=True)
        # Description
        add_textbox(slide, Inches(6.5), y + Inches(0.02), Inches(6), Inches(0.42),
                    desc, font_size=14, color=GRAY)

    # ═══════════════════════════════════════════════════════════════
    # SLIDE 8: Key Features & Differentiators
    # ═══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    pattern.seek(0)
    slide.shapes.add_picture(pattern, 0, 0, Emu(12192000), Emu(6858000))

    add_rounded_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT3)
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.7),
                "Ключевые особенности", font_size=36, color=WHITE, bold=True)

    features = [
        ("🌍", "Мультиязычность", "Русский и английский\nout-of-the-box.\nMiniLM мультиязычная\nмодель эмбеддингов.", ACCENT),
        ("🖥️", "Frontend", "React + TypeScript.\nИнтерактивный граф\nна D3.js. RAG-чат\nс историей запросов.", RGBColor(0x79, 0xC0, 0xFF)),
        ("🕸️", "Граф знаний", "Связи между\nперсонами, компаниями\nи локациями.\nОбогащение RAG-контекста.", ACCENT3),
        ("🛡️", "Отказоустойчивость", "3-уровневый fallback LLM.\nRetrieval-only режим.\nHealthcheck всех\nкомпонентов.", ORANGE),
        ("☁️", "Деплой", "Docker Compose локально.\nRailway + Vercel\nв облаке. Neo4j Aura\nи Qdrant Cloud.", RGBColor(0xFF, 0x60, 0x60)),
    ]

    for i, (icon, title, desc, color) in enumerate(features):
        x = Inches(0.3 + i * 2.6)
        add_rounded_rect(slide, x, Inches(1.5), Inches(2.4), Inches(4.0), BG_CARD)
        add_textbox(slide, x + Inches(0.15), Inches(1.65), Inches(2.1), Inches(0.5),
                    icon, font_size=36, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.15), Inches(2.25), Inches(2.1), Inches(0.45),
                    title, font_size=17, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_rounded_rect(slide, x + Inches(0.35), Inches(2.7), Inches(1.7), Pt(2), color)
        add_multiline(slide, x + Inches(0.2), Inches(2.9), Inches(2), Inches(2.2), [
            (desc, LIGHT),
        ], font_size=13, alignment=PP_ALIGN.CENTER, spacing=1.3)

    # Bottom: What's next
    add_rounded_rect(slide, Inches(0.5), Inches(5.8), Inches(12.3), Inches(1.2), BG_CARD)
    add_textbox(slide, Inches(0.9), Inches(5.9), Inches(3), Inches(0.4),
                "🚀  Ближайшие планы:", font_size=18, color=ACCENT, bold=True)
    add_multiline(slide, Inches(0.9), Inches(6.35), Inches(11.5), Inches(0.6), [
        ("Улучшение NER через LLM  •  Аналитика трендов  •  Пользовательские подписки  •  Streaming ответов  •  Celery workers", LIGHT),
    ], font_size=14)

    # ═══════════════════════════════════════════════════════════════
    # SLIDE 9: Thank you
    # ═══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    pattern.seek(0)
    slide.shapes.add_picture(pattern, 0, 0, Emu(12192000), Emu(6858000))

    # Graph decoration
    graph_img2 = create_node_graph_image(800, 600)
    slide.shapes.add_picture(graph_img2, Inches(7), Inches(1), Inches(6), Inches(4.5))

    add_rounded_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT)
    add_textbox(slide, Inches(1), Inches(1.5), Inches(6), Inches(1.2),
                "Спасибо за внимание!", font_size=44, color=WHITE, bold=True)
    add_textbox(slide, Inches(1), Inches(2.8), Inches(6), Inches(0.8),
                "IT News Platform", font_size=28, color=ACCENT)
    add_textbox(slide, Inches(1), Inches(3.5), Inches(6), Inches(0.8),
                "Интеллектуальная платформа агрегации IT-новостей\nс графом знаний и RAG",
                font_size=16, color=GRAY)

    # QR code for GitHub repo
    qr_img = make_qr("https://github.com/ArtemChik103/ITnews", size=240)
    slide.shapes.add_picture(qr_img, Inches(1), Inches(4.5), Inches(2), Inches(2))
    add_textbox(slide, Inches(1), Inches(6.5), Inches(2), Inches(0.4),
                "GitHub репозиторий", font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)

    # Contact info placeholder
    add_multiline(slide, Inches(3.5), Inches(4.8), Inches(4), Inches(1.5), [
        ("📧  [email]", LIGHT),
        ("💬  [telegram]", LIGHT),
        ("👥  [Название команды]", ACCENT),
    ], font_size=16, spacing=1.5)

    # Save
    output_path = os.path.join(OUT_DIR, "IT_News_Platform_Presentation_v2.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    return output_path


if __name__ == "__main__":
    build_presentation()
